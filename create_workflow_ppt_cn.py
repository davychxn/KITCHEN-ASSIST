"""
创建中文工作流程图PowerPoint演示文稿
锅具状态检测系统 - 使用最新工作流程更新
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path
from PIL import Image
import json


def add_rounded_rectangle(slide, left, top, width, height, text, fill_color, text_color=(255, 255, 255)):
    """添加带文本的圆角矩形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    
    # 设置填充颜色
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_color)
    
    # 设置边框
    shape.line.color.rgb = RGBColor(80, 80, 80)
    shape.line.width = Pt(1.5)
    
    # 添加文本
    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 格式化文本
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = RGBColor(*text_color)
            run.font.name = 'Microsoft YaHei'
    
    return shape


def add_arrow(slide, x1, y1, x2, y2, label=""):
    """在两点之间添加箭头连接器"""
    connector = slide.shapes.add_connector(
        1,  # 直线连接器
        x1, y1, x2, y2
    )
    connector.line.color.rgb = RGBColor(80, 80, 80)
    connector.line.width = Pt(2)
    
    # 添加箭头终点
    connector.line.end_arrow_type = 2
    
    # 如果提供了标签则添加
    if label:
        # 计算中点
        mid_x = (x1 + x2) / 2
        mid_y = (y1 + y2) / 2
        
        # 为标签添加文本框
        textbox = slide.shapes.add_textbox(
            mid_x - Inches(0.5), mid_y - Inches(0.2),
            Inches(1), Inches(0.3)
        )
        text_frame = textbox.text_frame
        text_frame.text = label
        text_frame.word_wrap = False
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(60, 60, 60)
                run.font.name = 'Microsoft YaHei'
    
    return connector


def create_workflow_presentation_cn():
    """创建完整的工作流程演示文稿"""
    
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 宽屏
    prs.slide_height = Inches(7.5)
    
    # ==================== 幻灯片1：标题页 ====================
    title_slide_layout = prs.slide_layouts[6]  # 空白布局
    slide1 = prs.slides.add_slide(title_slide_layout)
    
    # 背景颜色
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)  # 浅蓝色
    
    # 标题
    title_box = slide1.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(11.33), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = "锅具状态检测系统"
    for paragraph in title_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(48)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)
            run.font.name = 'Microsoft YaHei'
    
    # 副标题
    subtitle_box = slide1.shapes.add_textbox(
        Inches(1), Inches(3.8), Inches(11.33), Inches(0.6)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "MobileNet v2驱动的厨房安全检测"
    for paragraph in subtitle_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(24)
            run.font.color.rgb = RGBColor(100, 100, 100)
            run.font.name = 'Microsoft YaHei'
    
    # 添加统计数据
    stats_y = Inches(4.8)
    stats = [
        ("100% 训练与验证准确率", (76, 175, 80)),
        ("混合圆形+YOLO检测", (33, 150, 243)),
        ("MobileNet v2 - 350万参数", (255, 152, 0))
    ]
    
    x_start = Inches(2)
    spacing = Inches(3.2)
    
    for i, (stat, color) in enumerate(stats):
        stat_box = slide1.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_start + i * spacing, stats_y, Inches(2.8), Inches(0.6)
        )
        stat_box.fill.solid()
        stat_box.fill.fore_color.rgb = RGBColor(*color)
        stat_box.line.width = Pt(0)
        
        text_frame = stat_box.text_frame
        text_frame.text = stat
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.name = 'Microsoft YaHei'
    
    # ==================== 幻灯片2：整体工作流程 ====================
    slide2 = prs.slides.add_slide(title_slide_layout)
    
    # 背景
    background = slide2.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)
    
    # 标题
    title_box = slide2.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "系统工作流程概览"
    for paragraph in title_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)
            run.font.name = 'Microsoft YaHei'
    
    # 定义颜色
    color_input = (100, 149, 237)      # 矢车菊蓝
    color_detect = (255, 165, 0)       # 橙色
    color_classify = (50, 205, 50)     # 柠檬绿
    color_output = (220, 20, 60)       # 深红色
    
    # 方框尺寸 - 调整以获得更好的间距
    box_width = Inches(2.2)
    box_height = Inches(1)
    y_start = Inches(2)
    x_spacing = Inches(0.4)  # 减少方框之间的间距
    
    # 计算总宽度并将工作流程居中
    total_width = 4 * box_width + 3 * x_spacing
    x_start = (Inches(13.33) - total_width) / 2  # 水平居中
    
    # 步骤1：输入图像
    x1 = x_start
    box1 = add_rounded_rectangle(
        slide2, x1, y_start, box_width, box_height,
        "输入图像\\n(厨房相机)", color_input
    )
    
    # 步骤2：目标检测
    x2 = x1 + box_width + x_spacing
    box2 = add_rounded_rectangle(
        slide2, x2, y_start, box_width, box_height,
        "目标检测\\n(圆形+YOLO)", color_detect
    )
    
    # 步骤3：状态分类
    x3 = x2 + box_width + x_spacing
    box3 = add_rounded_rectangle(
        slide2, x3, y_start, box_width, box_height,
        "状态分类\\n(MobileNet v2)", color_classify
    )
    
    # 步骤4：输出
    x4 = x3 + box_width + x_spacing
    box4 = add_rounded_rectangle(
        slide2, x4, y_start, box_width, box_height,
        "输出结果\\n(带标记)", color_output
    )
    
    # 箭头
    add_arrow(slide2, 
              x1 + box_width, y_start + box_height/2,
              x2, y_start + box_height/2)
    
    add_arrow(slide2,
              x2 + box_width, y_start + box_height/2,
              x3, y_start + box_height/2)
    
    add_arrow(slide2,
              x3 + box_width, y_start + box_height/2,
              x4, y_start + box_height/2)
    
    # 在下方添加状态框
    y_states = Inches(4)
    state_width = Inches(1.5)
    state_height = Inches(0.6)
    
    # 状态颜色
    state_colors = {
        '正常': (0, 255, 0),
        '沸腾': (255, 255, 0),
        '冒烟': (128, 128, 128),
        '着火': (0, 0, 255)
    }
    
    states = list(state_colors.keys())
    x_state_start = Inches(2.5)
    state_spacing = Inches(2)
    
    for i, state in enumerate(states):
        x_state = x_state_start + i * state_spacing
        add_rounded_rectangle(
            slide2, x_state, y_states, state_width, state_height,
            state, state_colors[state], (0, 0, 0)
        )
    
    # 添加"检测状态："标签
    label_box = slide2.shapes.add_textbox(
        Inches(0.5), y_states, Inches(1.8), state_height
    )
    label_frame = label_box.text_frame
    label_frame.text = "检测状态："
    label_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in label_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 0, 0)
            run.font.name = 'Microsoft YaHei'
    
    # ==================== 幻灯片3：详细处理 ====================
    slide3 = prs.slides.add_slide(title_slide_layout)
    
    # 背景
    background = slide3.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)
    
    # 标题
    title_box = slide3.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "详细处理流程"
    for paragraph in title_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)
            run.font.name = 'Microsoft YaHei'
    
    # 左列 - 数据准备
    y_pos = Inches(1.3)
    x_left = Inches(0.8)
    box_w = Inches(2.5)
    box_h = Inches(0.8)
    spacing = Inches(0.9)
    
    # 数据准备部分
    section_box = slide3.shapes.add_textbox(x_left, y_pos, box_w, Inches(0.4))
    section_frame = section_box.text_frame
    section_frame.text = "数据准备"
    for paragraph in section_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)
            run.font.name = 'Microsoft YaHei'
    
    y_pos += Inches(0.5)
    
    # 数据准备步骤
    prep_steps = ["收集标记图像", "数据增强", "训练/验证分割"]
    for step in prep_steps:
        add_rounded_rectangle(slide3, x_left, y_pos, box_w, box_h,
                             step, (200, 200, 200), (0, 0, 0))
        y_pos += spacing
    
    # 中列 - 训练
    y_pos = Inches(1.3)
    x_mid = Inches(4.2)
    
    section_box = slide3.shapes.add_textbox(x_mid, y_pos, box_w, Inches(0.4))
    section_frame = section_box.text_frame
    section_frame.text = "模型训练"
    for paragraph in section_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)
            run.font.name = 'Microsoft YaHei'
    
    y_pos += Inches(0.5)
    
    add_rounded_rectangle(slide3, x_mid, y_pos, box_w, box_h,
                         "训练分类器\\n(MobileNet v2)", (144, 238, 144), (0, 0, 0))
    y_pos += spacing
    
    add_rounded_rectangle(slide3, x_mid, y_pos, box_w, box_h,
                         "评估性能", (255, 215, 0), (0, 0, 0))
    y_pos += spacing
    
    add_rounded_rectangle(slide3, x_mid, y_pos, box_w, box_h,
                         "保存最佳模型", (135, 206, 250), (0, 0, 0))
    
    # 右列 - 预测
    y_pos = Inches(1.3)
    x_right = Inches(7.6)
    
    section_box = slide3.shapes.add_textbox(x_right, y_pos, box_w, Inches(0.4))
    section_frame = section_box.text_frame
    section_frame.text = "生产预测"
    for paragraph in section_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)
            run.font.name = 'Microsoft YaHei'
    
    y_pos += Inches(0.5)
    
    pred_steps = ["圆形检测", "YOLO后备", "状态分类"]
    for step in pred_steps:
        add_rounded_rectangle(slide3, x_right, y_pos, box_w, box_h,
                             step, (255, 182, 193), (0, 0, 0))
        y_pos += spacing
    
    # 最远右列 - 输出
    y_pos = Inches(1.3)
    x_far_right = Inches(11)
    
    section_box = slide3.shapes.add_textbox(x_far_right, y_pos, box_w, Inches(0.4))
    section_frame = section_box.text_frame
    section_frame.text = "结果输出"
    for paragraph in section_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)
            run.font.name = 'Microsoft YaHei'
    
    y_pos += Inches(0.5)
    
    output_steps = ["标记图像", "JSON报告", "性能指标"]
    for step in output_steps:
        add_rounded_rectangle(slide3, x_far_right, y_pos, box_w, box_h,
                             step, (221, 160, 221), (0, 0, 0))
        y_pos += spacing
    
    # ==================== 幻灯片4：主要特点 ====================
    slide4 = prs.slides.add_slide(title_slide_layout)
    
    # 背景
    background = slide4.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)
    
    # 标题
    title_box = slide4.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "主要特点"
    for paragraph in title_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)
            run.font.name = 'Microsoft YaHei'
    
    # 特点
    features = [
        ("🎯 混合检测", "圆形检测 + YOLO v8后备"),
        ("🧠 轻量级模型", "MobileNet v2（350万参数）"),
        ("📊 高准确率", "训练与验证100%"),
        ("🖼️ 精准线框", "紧密贴合圆形锅具"),
        ("🎨 颜色优化", "保留关键颜色特征"),
        ("⚡ 快速推理", "针对边缘设备优化")
    ]
    
    y_start = Inches(1.5)
    x_left = Inches(1.5)
    x_right = Inches(7)
    feature_height = Inches(0.8)
    feature_spacing = Inches(1)
    
    for i, (title, desc) in enumerate(features):
        row = i // 2
        col = i % 2
        x = x_left if col == 0 else x_right
        y = y_start + row * feature_spacing
        
        # 特点框
        feature_box = slide4.shapes.add_textbox(
            x, y, Inches(5), feature_height
        )
        text_frame = feature_box.text_frame
        text_frame.text = f"{title}\\n{desc}"
        
        for paragraph in text_frame.paragraphs:
            if paragraph.text.startswith("🎯") or paragraph.text.startswith("🧠") or \
               paragraph.text.startswith("📊") or paragraph.text.startswith("🖼️") or \
               paragraph.text.startswith("🎨") or paragraph.text.startswith("⚡"):
                for run in paragraph.runs:
                    run.font.size = Pt(18)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0, 51, 102)
                    run.font.name = 'Microsoft YaHei'
            else:
                for run in paragraph.runs:
                    run.font.size = Pt(14)
                    run.font.color.rgb = RGBColor(60, 60, 60)
                    run.font.name = 'Microsoft YaHei'
    
    # ==================== 幻灯片5：性能结果 ====================
    slide5 = prs.slides.add_slide(title_slide_layout)
    
    # 背景
    background = slide5.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)
    
    # 标题
    title_box = slide5.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "性能结果"
    for paragraph in title_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)
            run.font.name = 'Microsoft YaHei'
    
    # 性能指标
    y_pos = Inches(1.5)
    
    # 训练准确率
    metric_box = slide5.shapes.add_textbox(
        Inches(2), y_pos, Inches(9), Inches(0.8)
    )
    metric_frame = metric_box.text_frame
    metric_frame.text = "训练准确率：100%（40/40张图像）"
    for paragraph in metric_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(24)
            run.font.bold = True
            run.font.color.rgb = RGBColor(76, 175, 80)
            run.font.name = 'Microsoft YaHei'
    
    y_pos += Inches(1)
    
    # 验证准确率
    metric_box = slide5.shapes.add_textbox(
        Inches(2), y_pos, Inches(9), Inches(0.8)
    )
    metric_frame = metric_box.text_frame
    metric_frame.text = "验证准确率：100%（4/4张图像）"
    for paragraph in metric_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(24)
            run.font.bold = True
            run.font.color.rgb = RGBColor(33, 150, 243)
            run.font.name = 'Microsoft YaHei'
    
    y_pos += Inches(1.2)
    
    # 模型信息
    info_text = [
        "模型：MobileNet v2",
        "参数：350万（比ResNet18少68%）",
        "训练：200轮，批次大小4",
        "检测：混合圆形+YOLO方法"
    ]
    
    for line in info_text:
        info_box = slide5.shapes.add_textbox(
            Inches(3), y_pos, Inches(7), Inches(0.4)
        )
        info_frame = info_box.text_frame
        info_frame.text = line
        for paragraph in info_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(60, 60, 60)
                run.font.name = 'Microsoft YaHei'
        y_pos += Inches(0.5)
    
    # 添加预测结果（如果有）
    predictions_file = Path('./veri_results_marked/predictions.json')
    if predictions_file.exists():
        try:
            with open(predictions_file, 'r') as f:
                predictions = json.load(f)
            
            if predictions:
                y_pos += Inches(0.2)
                pred_title = slide5.shapes.add_textbox(
                    Inches(1), y_pos, Inches(11), Inches(0.4)
                )
                pred_frame = pred_title.text_frame
                pred_frame.text = "最新预测结果："
                for paragraph in pred_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(18)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0, 51, 102)
                        run.font.name = 'Microsoft YaHei'
                
                y_pos += Inches(0.5)
                
                # 显示前4个预测
                for pred in predictions[:4]:
                    filename = pred['filename']
                    pred_state = pred['predicted_state']
                    confidence = pred['confidence']
                    
                    state_map = {
                        'boiling': '沸腾',
                        'normal': '正常', 
                        'on_fire': '着火',
                        'smoking': '冒烟'
                    }
                    pred_state_cn = state_map.get(pred_state, pred_state)
                    
                    pred_box = slide5.shapes.add_textbox(
                        Inches(2), y_pos, Inches(9), Inches(0.3)
                    )
                    pred_frame = pred_box.text_frame
                    pred_frame.text = f"{filename}: {pred_state_cn} (置信度: {confidence:.3f})"
                    for paragraph in pred_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(12)
                            run.font.color.rgb = RGBColor(60, 60, 60)
                            run.font.name = 'Microsoft YaHei'
                    
                    y_pos += Inches(0.4)
                    if y_pos > Inches(6.5):
                        break
        except Exception as e:
            print(f"无法加载预测结果：{e}")
    
    # ==================== 幻灯片6：标记图像 ====================
    slide6 = prs.slides.add_slide(title_slide_layout)
    
    # 背景
    background = slide6.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)
    
    # 标题
    title_box = slide6.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "标记检测结果"
    for paragraph in title_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)
            run.font.name = 'Microsoft YaHei'
    
    # 添加标记图像（如果有）- 所有图像缩放到相似大小
    marked_dir = Path('./veri_results_marked')
    if marked_dir.exists():
        marked_images = sorted(list(marked_dir.glob('*_marked.jpg')))[:4]  # 前4张图像
        
        if marked_images:
            # 第一遍：找到尺寸以将所有图像标准化为相似大小
            image_dims = []
            for img_path in marked_images:
                try:
                    with Image.open(img_path) as img:
                        image_dims.append((img.size[0], img.size[1]))
                except:
                    image_dims.append((800, 600))
            
            # 使用平均纵横比确定统一大小
            avg_aspect = sum(w/h for w, h in image_dims) / len(image_dims)
            
            # 根据可用空间设置统一目标大小
            if avg_aspect > 1.5:  # 更宽的图像
                target_width = Inches(5.8)
                target_height = Inches(target_width.inches / avg_aspect)
            else:  # 更方形或更高的图像
                target_height = Inches(2.6)
                target_width = Inches(target_height.inches * avg_aspect)
            
            # 以统一大小添加2x2网格中的图像
            x_positions = [Inches(0.8), Inches(7.2)]
            y_positions = [Inches(1.2), Inches(4.2)]
            
            for idx, img_path in enumerate(marked_images):
                if idx >= 4:
                    break
                row = idx // 2
                col = idx % 2
                x = x_positions[col]
                y = y_positions[row]
                
                try:
                    # 以统一大小添加图像
                    pic = slide6.shapes.add_picture(
                        str(img_path), x, y, width=target_width, height=target_height
                    )
                    
                    # 在图像下方添加标题
                    caption_box = slide6.shapes.add_textbox(
                        x, y + target_height + Inches(0.05), target_width, Inches(0.3)
                    )
                    caption_frame = caption_box.text_frame
                    caption_text = img_path.stem.replace('_marked', '').replace('_', ' ')
                    caption_frame.text = caption_text[:40] + ('...' if len(caption_text) > 40 else '')
                    for paragraph in caption_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                        for run in paragraph.runs:
                            run.font.size = Pt(10)
                            run.font.color.rgb = RGBColor(80, 80, 80)
                            run.font.name = 'Microsoft YaHei'
                except Exception as e:
                    print(f"无法添加图像 {img_path.name}：{e}")
        else:
            # 无图像消息
            msg_box = slide6.shapes.add_textbox(
                Inches(2), Inches(3), Inches(9), Inches(1)
            )
            msg_frame = msg_box.text_frame
            msg_frame.text = "运行 predict_veri.py 生成标记图像"
            msg_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in msg_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(24)
                    run.font.color.rgb = RGBColor(150, 150, 150)
                    run.font.name = 'Microsoft YaHei'
    
    # 保存演示文稿
    prs.save('Kitchen_Assist_Workflow_CN.pptx')
    print("✓ 中文PowerPoint演示文稿已创建：Kitchen_Assist_Workflow_CN.pptx")
    print(f"  - 总幻灯片数：{len(prs.slides)}")
    if predictions_file.exists():
        with open(predictions_file, 'r') as f:
            predictions = json.load(f)
        print(f"  - 包含 {len(predictions)} 个预测结果")
    if marked_dir.exists() and marked_images:
        print(f"  - 包含 {len(marked_images)} 张标记图像")


if __name__ == "__main__":
    create_workflow_presentation_cn()
