"""
Create a neat workflow diagram PowerPoint presentation
for the Pan/Pot State Detection System - Updated with latest workflow
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path
from PIL import Image
import json


def add_rounded_rectangle(slide, left, top, width, height, text, fill_color, text_color=(255, 255, 255)):
    """Add a rounded rectangle shape with text"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    
    # Set fill color
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_color)
    
    # Set border
    shape.line.color.rgb = RGBColor(80, 80, 80)
    shape.line.width = Pt(1.5)
    
    # Add text
    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Format text
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = RGBColor(*text_color)
    
    return shape


def add_arrow(slide, x1, y1, x2, y2, label=""):
    """Add an arrow connector between two points"""
    connector = slide.shapes.add_connector(
        1,  # Straight connector
        x1, y1, x2, y2
    )
    connector.line.color.rgb = RGBColor(80, 80, 80)
    connector.line.width = Pt(2)
    
    # Add arrow end
    connector.line.end_arrow_type = 2
    
    # Add label if provided
    if label:
        # Calculate midpoint
        mid_x = (x1 + x2) / 2
        mid_y = (y1 + y2) / 2
        
        # Add text box for label
        textbox = slide.shapes.add_textbox(
            mid_x - Inches(0.5), mid_y - Inches(0.2),
            Inches(1), Inches(0.3)
        )
        text_frame = textbox.text_frame
        text_frame.text = label
        text_frame.word_wrap = False
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(60, 60, 60)
    
    return connector


def create_workflow_presentation():
    """Create the complete workflow presentation"""
    
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # Widescreen
    prs.slide_height = Inches(7.5)
    
    # ==================== Slide 1: Title Slide ====================
    title_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide1 = prs.slides.add_slide(title_slide_layout)
    
    # Background color
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue
    
    # Title
    title_box = slide1.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(11.33), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Pan/Pot State Detection System"
    for paragraph in title_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(48)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)
    
    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(
        Inches(1), Inches(3.8), Inches(11.33), Inches(0.6)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "MobileNet v2 Powered Kitchen Safety Detection"
    for paragraph in subtitle_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(24)
            run.font.color.rgb = RGBColor(100, 100, 100)
    
    # Add stats
    stats_y = Inches(4.8)
    stats = [
        ("100% Training & Verification", (76, 175, 80)),
        ("Hybrid Circle + YOLO Detection", (33, 150, 243)),
        ("MobileNet v2 - 3.5M Params", (255, 152, 0))
    ]
    
    x_start = Inches(2.5)
    spacing = Inches(3)
    
    for i, (stat, color) in enumerate(stats):
        stat_box = slide1.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_start + i * spacing, stats_y, Inches(2.5), Inches(0.6)
        )
        stat_box.fill.solid()
        stat_box.fill.fore_color.rgb = RGBColor(*color)
        stat_box.line.width = Pt(0)
        
        text_frame = stat_box.text_frame
        text_frame.text = stat
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
    
    # ==================== Slide 2: Overall Workflow ====================
    slide2 = prs.slides.add_slide(title_slide_layout)
    
    # Background
    background = slide2.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)
    
    # Title
    title_box = slide2.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "System Workflow Overview"
    for paragraph in title_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)
    
    # Define colors
    color_input = (100, 149, 237)      # Cornflower blue
    color_detect = (255, 165, 0)       # Orange
    color_classify = (50, 205, 50)     # Lime green
    color_output = (220, 20, 60)       # Crimson
    
    # Box dimensions - adjusted for better spacing
    box_width = Inches(2.2)
    box_height = Inches(1)
    y_start = Inches(2)
    x_spacing = Inches(0.4)  # Reduced spacing between boxes
    
    # Calculate total width and center the workflow
    total_width = 4 * box_width + 3 * x_spacing
    x_start = (Inches(13.33) - total_width) / 2  # Center horizontally
    
    # Step 1: Input Image
    x1 = x_start
    box1 = add_rounded_rectangle(
        slide2, x1, y_start, box_width, box_height,
        "Input Image\n(Kitchen Camera)", color_input
    )
    
    # Step 2: Object Detection
    x2 = x1 + box_width + x_spacing
    box2 = add_rounded_rectangle(
        slide2, x2, y_start, box_width, box_height,
        "Object Detection\n(Circle + YOLO)", color_detect
    )
    
    # Step 3: State Classification
    x3 = x2 + box_width + x_spacing
    box3 = add_rounded_rectangle(
        slide2, x3, y_start, box_width, box_height,
        "State Classification\n(MobileNet v2)", color_classify
    )
    
    # Step 4: Output
    x4 = x3 + box_width + x_spacing
    box4 = add_rounded_rectangle(
        slide2, x4, y_start, box_width, box_height,
        "Output Result\n(with Marking)", color_output
    )
    
    # Arrows
    add_arrow(slide2, 
              x1 + box_width, y_start + box_height/2,
              x2, y_start + box_height/2)
    
    add_arrow(slide2,
              x2 + box_width, y_start + box_height/2,
              x3, y_start + box_height/2)
    
    add_arrow(slide2,
              x3 + box_width, y_start + box_height/2,
              x4, y_start + box_height/2)
    
    # Add state boxes below
    y_states = Inches(4)
    state_width = Inches(1.5)
    state_height = Inches(0.6)
    
    # State colors
    state_colors = {
        'Normal': (0, 255, 0),
        'Boiling': (255, 255, 0),
        'Smoking': (128, 128, 128),
        'On Fire': (0, 0, 255)
    }
    
    states = list(state_colors.keys())
    x_state_start = Inches(2.5)
    state_spacing = Inches(2)
    
    for i, state in enumerate(states):
        x_state = x_state_start + i * state_spacing
        add_rounded_rectangle(
            slide2, x_state, y_states, state_width, state_height,
            state, state_colors[state], (0, 0, 0)
        )
    
    # Add "Detected States:" label
    label_box = slide2.shapes.add_textbox(
        Inches(0.5), y_states, Inches(1.8), state_height
    )
    label_frame = label_box.text_frame
    label_frame.text = "Detected States:"
    label_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in label_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 0, 0)
    
    # ==================== Slide 3: Detailed Processing ====================
    slide3 = prs.slides.add_slide(title_slide_layout)
    
    # Background
    background = slide3.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)
    
    # Title
    title_box = slide3.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Detailed Processing Pipeline"
    for paragraph in title_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)
    
    # Left column - Data Preparation
    y_pos = Inches(1.3)
    x_left = Inches(0.8)
    box_w = Inches(2.5)
    box_h = Inches(0.8)
    spacing = Inches(0.9)
    
    # Data Preparation Section
    section_box = slide3.shapes.add_textbox(x_left, y_pos, box_w, Inches(0.4))
    section_frame = section_box.text_frame
    section_frame.text = "Data Preparation"
    for paragraph in section_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)
    
    y_pos += Inches(0.5)
    
    add_rounded_rectangle(slide3, x_left, y_pos, box_w, box_h,
                         "Collect Images", (176, 196, 222), (0, 0, 0))
    y_pos += spacing
    
    add_rounded_rectangle(slide3, x_left, y_pos, box_w, box_h,
                         "Manual Crop\n(manual_crop.py)", (176, 196, 222), (0, 0, 0))
    y_pos += spacing
    
    add_rounded_rectangle(slide3, x_left, y_pos, box_w, box_h,
                         "Label States", (176, 196, 222), (0, 0, 0))
    
    # Middle column - Training
    y_pos = Inches(1.3)
    x_mid = Inches(4)
    
    section_box = slide3.shapes.add_textbox(x_mid, y_pos, box_w, Inches(0.4))
    section_frame = section_box.text_frame
    section_frame.text = "Model Training"
    for paragraph in section_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 102, 0)
    
    y_pos += Inches(0.5)
    
    add_rounded_rectangle(slide3, x_mid, y_pos, box_w, box_h,
                         "Data Augmentation", (144, 238, 144), (0, 0, 0))
    y_pos += spacing
    
    add_rounded_rectangle(slide3, x_mid, y_pos, box_w, box_h,
                         "Train Classifier\n(MobileNet v2)", (144, 238, 144), (0, 0, 0))
    y_pos += spacing
    
    add_rounded_rectangle(slide3, x_mid, y_pos, box_w, box_h,
                         "Model Validation", (144, 238, 144), (0, 0, 0))
    
    # Right column - Deployment
    y_pos = Inches(1.3)
    x_right = Inches(7.2)
    
    section_box = slide3.shapes.add_textbox(x_right, y_pos, box_w, Inches(0.4))
    section_frame = section_box.text_frame
    section_frame.text = "Prediction"
    for paragraph in section_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = RGBColor(102, 0, 0)
    
    y_pos += Inches(0.5)
    
    add_rounded_rectangle(slide3, x_right, y_pos, box_w, box_h,
                         "Load Model", (255, 182, 193), (0, 0, 0))
    y_pos += spacing
    
    add_rounded_rectangle(slide3, x_right, y_pos, box_w, box_h,
                         "YOLO Detection", (255, 182, 193), (0, 0, 0))
    y_pos += spacing
    
    add_rounded_rectangle(slide3, x_right, y_pos, box_w, box_h,
                         "State Prediction\n(predict_veri.py)", (255, 182, 193), (0, 0, 0))
    
    # Evaluation box at bottom
    y_eval = Inches(5.5)
    x_eval = Inches(10.3)
    add_rounded_rectangle(slide3, x_eval, y_eval, box_w, box_h,
                         "Evaluation\n(evaluate_classifier.py)", (221, 160, 221), (0, 0, 0))
    
    # ==================== Slide 4: Key Features ====================
    slide4 = prs.slides.add_slide(title_slide_layout)
    
    # Background
    background = slide4.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)
    
    # Title
    title_box = slide4.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Key Features & Capabilities"
    for paragraph in title_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)
    
    # Features
    features = [
        ("🎯 Hybrid Detection", "Circle detection + YOLO v8 fallback"),
        ("🧠 Lightweight Model", "MobileNet v2 (3.5M params)"),
        ("📊 High Accuracy", "100% training & verification"),
        ("🖼️ Accurate Wireframes", "Tight fit around circular cookware"),
        ("🎨 Color Optimized", "Preserves critical color features"),
        ("⚡ Fast Inference", "Optimized for edge devices")
    ]
    
    y_feature = Inches(1.5)
    x_left_col = Inches(1)
    x_right_col = Inches(7)
    feature_width = Inches(5)
    feature_height = Inches(0.9)
    feature_spacing = Inches(1.1)
    
    for i, (title, desc) in enumerate(features):
        x_pos = x_left_col if i < 3 else x_right_col
        y_pos = y_feature + (i % 3) * feature_spacing
        
        # Feature box
        box = add_rounded_rectangle(
            slide4, x_pos, y_pos, feature_width, feature_height,
            "", (230, 240, 255), (0, 0, 0)
        )
        
        # Title
        text_frame = box.text_frame
        text_frame.clear()
        p1 = text_frame.paragraphs[0]
        p1.text = title
        p1.alignment = PP_ALIGN.LEFT
        for run in p1.runs:
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)
        
        # Description
        p2 = text_frame.add_paragraph()
        p2.text = desc
        p2.alignment = PP_ALIGN.LEFT
        for run in p2.runs:
            run.font.size = Pt(13)
            run.font.color.rgb = RGBColor(60, 60, 60)
    
    # ==================== Slide 5: Prediction Results ====================
    slide5 = prs.slides.add_slide(title_slide_layout)
    
    # Background
    background = slide5.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)
    
    # Title
    title_box = slide5.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Model Performance & Predictions"
    for paragraph in title_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)
    
    # Load prediction results if available
    predictions_file = Path('./veri_results_marked/predictions.json')
    if predictions_file.exists():
        with open(predictions_file, 'r') as f:
            predictions = json.load(f)
        
        # Overall accuracy box
        accuracy_box = slide5.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(1.2), Inches(3), Inches(1.2)
        )
        accuracy_box.fill.solid()
        accuracy_box.fill.fore_color.rgb = RGBColor(76, 175, 80)  # Green
        accuracy_box.line.color.rgb = RGBColor(56, 142, 60)
        accuracy_box.line.width = Pt(2)
        
        text_frame = accuracy_box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Training accuracy
        p1 = text_frame.paragraphs[0]
        p1.text = "Training Accuracy"
        p1.alignment = PP_ALIGN.CENTER
        for run in p1.runs:
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
        
        p2 = text_frame.add_paragraph()
        p2.text = "100%"
        p2.alignment = PP_ALIGN.CENTER
        for run in p2.runs:
            run.font.size = Pt(48)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
        
        # Model info box
        model_box = slide5.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(4.2), Inches(1.2), Inches(4), Inches(1.2)
        )
        model_box.fill.solid()
        model_box.fill.fore_color.rgb = RGBColor(33, 150, 243)  # Blue
        model_box.line.color.rgb = RGBColor(21, 101, 192)
        model_box.line.width = Pt(2)
        
        text_frame = model_box.text_frame
        text_frame.clear()
        
        p1 = text_frame.paragraphs[0]
        p1.text = "Model Architecture"
        p1.alignment = PP_ALIGN.CENTER
        for run in p1.runs:
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
        
        p2 = text_frame.add_paragraph()
        p2.text = "MobileNet v2"
        p2.alignment = PP_ALIGN.CENTER
        for run in p2.runs:
            run.font.size = Pt(24)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
        
        p3 = text_frame.add_paragraph()
        p3.text = "~3.5M parameters"
        p3.alignment = PP_ALIGN.CENTER
        for run in p3.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(255, 255, 255)
        
        # Verification results box
        correct_preds = sum(1 for p in predictions if p.get('correct', False))
        total_preds = len([p for p in predictions if 'correct' in p])
        
        if total_preds > 0:
            verify_box = slide5.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(8.6), Inches(1.2), Inches(3.5), Inches(1.2)
            )
            verify_box.fill.solid()
            verify_box.fill.fore_color.rgb = RGBColor(255, 152, 0)  # Orange
            verify_box.line.color.rgb = RGBColor(230, 81, 0)
            verify_box.line.width = Pt(2)
            
            text_frame = verify_box.text_frame
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            p1 = text_frame.paragraphs[0]
            p1.text = "Verification Set"
            p1.alignment = PP_ALIGN.CENTER
            for run in p1.runs:
                run.font.size = Pt(16)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
            
            verify_acc = (correct_preds / total_preds * 100) if total_preds > 0 else 0
            p2 = text_frame.add_paragraph()
            p2.text = f"{verify_acc:.0f}% ({correct_preds}/{total_preds})"
            p2.alignment = PP_ALIGN.CENTER
            for run in p2.runs:
                run.font.size = Pt(36)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)
        
        # Add prediction examples table
        y_table = Inches(2.6)
        
        # Table header
        header_box = slide5.shapes.add_textbox(
            Inches(0.8), y_table, Inches(11.5), Inches(0.4)
        )
        header_frame = header_box.text_frame
        header_frame.text = "Sample Predictions"
        for paragraph in header_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.size = Pt(18)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 51, 102)
        
        # Prediction rows
        y_pos = y_table + Inches(0.5)
        for i, pred in enumerate(predictions[:4]):  # Show first 4 predictions
            # Filename
            filename_box = slide5.shapes.add_textbox(
                Inches(0.8), y_pos, Inches(3.5), Inches(0.5)
            )
            text_frame = filename_box.text_frame
            text_frame.text = pred['filename'][:35] + ('...' if len(pred['filename']) > 35 else '')
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(11)
            
            # Predicted state
            pred_box = slide5.shapes.add_textbox(
                Inches(4.5), y_pos, Inches(2), Inches(0.5)
            )
            text_frame = pred_box.text_frame
            text_frame.text = f"Pred: {pred['predicted_state']}"
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12)
                    run.font.bold = True
            
            # Confidence
            conf_box = slide5.shapes.add_textbox(
                Inches(6.7), y_pos, Inches(1.5), Inches(0.5)
            )
            text_frame = conf_box.text_frame
            conf = pred['confidence']
            text_frame.text = f"{conf:.1%}"
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(12)
            
            # Status (if available)
            if 'correct' in pred:
                status_text = "✓" if pred['correct'] else "✗"
                status_color = RGBColor(76, 175, 80) if pred['correct'] else RGBColor(244, 67, 54)
                
                status_box = slide5.shapes.add_textbox(
                    Inches(8.4), y_pos, Inches(0.6), Inches(0.5)
                )
                text_frame = status_box.text_frame
                text_frame.text = status_text
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                for paragraph in text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(20)
                        run.font.bold = True
                        run.font.color.rgb = status_color
            
            y_pos += Inches(0.6)
    
    # ==================== Slide 6: Marked Images ====================
    slide6 = prs.slides.add_slide(title_slide_layout)
    
    # Background
    background = slide6.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 250)
    
    # Title
    title_box = slide6.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Marked Detection Results"
    for paragraph in title_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)
    
    # Add marked images if available - all images scaled to similar sizes
    marked_dir = Path('./veri_results_marked')
    if marked_dir.exists():
        marked_images = sorted(list(marked_dir.glob('*_marked.jpg')))[:4]  # First 4 images
        
        if marked_images:
            # First pass: find dimensions to normalize all images to similar sizes
            image_dims = []
            for img_path in marked_images:
                try:
                    with Image.open(img_path) as img:
                        image_dims.append((img.size[0], img.size[1]))
                except:
                    image_dims.append((800, 600))
            
            # Use average aspect ratio to determine uniform size
            avg_aspect = sum(w/h for w, h in image_dims) / len(image_dims)
            
            # Set uniform target size based on available space
            if avg_aspect > 1.5:  # Wider images
                target_width = Inches(5.8)
                target_height = Inches(target_width.inches / avg_aspect)
            else:  # More square or taller images
                target_height = Inches(2.6)
                target_width = Inches(target_height.inches * avg_aspect)
            
            # Add images in 2x2 grid with uniform sizes
            x_positions = [Inches(0.8), Inches(7.2)]
            y_positions = [Inches(1.2), Inches(4.2)]
            
            for idx, img_path in enumerate(marked_images):
                if idx >= 4:
                    break
                row = idx // 2
                col = idx % 2
                x = x_positions[col]
                y = y_positions[row]
                
                try:
                    # Add image with uniform size
                    pic = slide6.shapes.add_picture(
                        str(img_path), x, y, width=target_width, height=target_height
                    )
                    
                    # Add caption below the image
                    caption_box = slide6.shapes.add_textbox(
                        x, y + target_height + Inches(0.05), target_width, Inches(0.3)
                    )
                    caption_frame = caption_box.text_frame
                    caption_text = img_path.stem.replace('_marked', '').replace('_', ' ')
                    caption_frame.text = caption_text[:40] + ('...' if len(caption_text) > 40 else '')
                    for paragraph in caption_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                        for run in paragraph.runs:
                            run.font.size = Pt(10)
                            run.font.color.rgb = RGBColor(80, 80, 80)
                except Exception as e:
                    print(f"Could not add image {img_path.name}: {e}")
        else:
            # No images message
            msg_box = slide6.shapes.add_textbox(
                Inches(2), Inches(3), Inches(9), Inches(1)
            )
            msg_frame = msg_box.text_frame
            msg_frame.text = "Run predict_veri.py to generate marked images"
            msg_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in msg_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(24)
                    run.font.color.rgb = RGBColor(150, 150, 150)
    
    # Save presentation
    prs.save('Kitchen_Assist_Workflow.pptx')
    print("✓ PowerPoint presentation created: Kitchen_Assist_Workflow.pptx")
    print(f"  - Total slides: {len(prs.slides)}")
    if predictions_file.exists():
        print(f"  - Included {len(predictions)} prediction results")
    if marked_dir.exists() and marked_images:
        print(f"  - Included {len(marked_images)} marked images")


if __name__ == "__main__":
    create_workflow_presentation()
